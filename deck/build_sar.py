# -*- coding: utf-8 -*-
"""Build the ChromaSAR (SIH1733) idea deck, with the architecture diagram and
speaker notes for the team.

    python deck/build_sar.py

Template rules honoured:
  * all 7 slides kept - the coordinator's email says no slides may be added or deleted,
    which is stricter than the template's own "max 6" note, so we obey the stricter one
  * the "idea details pointers" (slide titles, and slide 2's underlined Proposed
    Solution header) are preserved
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deckkit import ARIAL, card, fill, fit_fill, fit_group, place, shp, stat  # noqa: E402
import content_sar as C                                                       # noqa: E402

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
TEMPLATE = os.path.join(ROOT, "SIH2026-IDEA-Presentation-Format.pptx")
DIAGRAM = os.path.join(HERE, "architecture_slide.png")
OUT = os.path.join(ROOT, "SIH2026-DeltaForce-SIH1733-SAR-Colorization.pptx")

TOP, BOT, L, R = 1.30, 6.82, 0.50, 12.83
W = R - L
P, A, BODY, MUTED = (C.PAL["primary"], C.PAL["accent"], C.PAL["body"], C.PAL["muted"])


def drop(shape):
    shape._element.getparent().remove(shape._element)


def set_title(slide, text=None, size=None):
    t = shp(slide, [s.name for s in slide.shapes if s.name.startswith("Title")][0])
    place(t, 1.85, 0.02, 8.70, 1.10)
    t.text_frame.word_wrap = True
    t.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    if text is not None:
        fill(t.text_frame, [{"runs": [(text, True)], "size": size, "color": P,
                             "line": 92}])
    else:
        for p in t.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT


def set_team(slide):
    for s in slide.shapes:
        if s.name.startswith("Oval"):
            s.text_frame.word_wrap = True
            s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            s.text_frame.margin_left = Inches(0.03)
            s.text_frame.margin_right = Inches(0.03)
            fill(s.text_frame, [{"runs": [(C.TEAM, True)], "size": 11, "color": P,
                                 "align": "c", "line": 95}])
            return


def col_card(slide, l, t, w, h, key="card"):
    c = card(slide, l, t, w, h, C.PAL[key],
             C.PAL["cardline"] if key == "card" else C.PAL["card2line"])
    tf = c.text_frame
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.13)
    tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return c


def notes(slide, text):
    """Speaker notes - what the presenter actually says, plus expected questions."""
    slide.notes_slide.notes_text_frame.text = text


def main():
    prs = Presentation(TEMPLATE)
    S = prs.slides

    # ---------------- 1 title ------------------------------------------
    s1 = S[0]
    tb = shp(s1, "TextBox 9")
    place(tb, 0.36, 2.05, 5.72, 4.65)
    rows = [("Problem Statement ID - ", C.PS_ID),
            ("Problem Statement Title - ", C.PS_TITLE),
            ("Theme - ", C.THEME),
            ("PS Category - ", C.CATEGORY),
            ("Team ID - ", C.TEAM_ID),
            ("Team Name - ", C.TEAM)]
    fit_fill(tb, [{"runs": [(k, True), (v, False)], "size": 15, "color": BODY,
                   "bullet": "•", "after": 9, "line": 108} for k, v in rows],
             cap=1.30)
    notes(s1, C.NOTES[1])

    # ---------------- 2 idea -------------------------------------------
    s2 = S[1]
    set_title(s2, C.IDEA_TITLE, 19)
    set_team(s2)
    body = shp(s2, "TextBox 8")
    place(body, L, TOP, 7.72, BOT - TOP)
    blk = [{"runs": [("Proposed Solution (Describe your Idea/Solution/Prototype)", True)],
            "size": 15, "color": P, "underline": True, "after": 10}]
    for lab, txt in C.S2_LEFT:
        blk.append({"runs": [(lab, True), (txt, False)], "size": 12.5, "color": BODY,
                    "bullet": "▪", "after": 9, "line": 104})
    fit_fill(body, blk, cap=1.45)

    c = col_card(s2, 8.55, TOP, R - 8.55, BOT - TOP, "card2")
    blk = [{"runs": [(C.S2_RIGHT_HDR, True)], "size": 12.5, "color": A, "after": 10}]
    for h, t in C.S2_RIGHT:
        blk.append({"runs": [(h, True)], "size": 11.5, "color": P, "after": 2,
                    "line": 102})
        blk.append({"runs": [(t, False)], "size": 10.5, "color": MUTED, "after": 10,
                    "line": 104})
    fit_fill(c, blk, cap=1.45)
    notes(s2, C.NOTES[2])

    # ---------------- 3 technical (diagram + two columns) ---------------
    s3 = S[2]
    set_title(s3)
    set_team(s3)
    drop(shp(s3, "TextBox 8"))
    dw = 11.0
    s3.shapes.add_picture(DIAGRAM, Inches(L + (W - dw) / 2), Inches(TOP),
                          width=Inches(dw), height=Inches(dw / 4.0))
    ct = TOP + dw / 4.0 + 0.22
    ch = BOT - ct
    cw = (W - 0.45) / 2
    pairs = []
    c = col_card(s3, L, ct, cw, ch)
    blk = [{"runs": [(C.S3_LEFT_HDR, True)], "size": 11.5, "color": A, "after": 8}]
    for lab, txt in C.S3_LEFT:
        blk.append({"runs": [(lab, True), (txt, False)], "size": 10.5, "color": BODY,
                    "bullet": "▪", "after": 5, "line": 104})
    pairs.append((c, blk))
    c = col_card(s3, L + cw + 0.45, ct, cw, ch)
    blk = [{"runs": [(C.S3_RIGHT_HDR, True)], "size": 11.5, "color": A, "after": 8}]
    for lab, txt in C.S3_RIGHT:
        blk.append({"runs": [(lab, True), (" — ", True), (txt, False)], "size": 10.5,
                    "color": BODY, "bullet": "▪", "after": 5, "line": 104})
    pairs.append((c, blk))
    fit_group(pairs, cap=1.4)
    notes(s3, C.NOTES[3])

    # ---------------- 4 feasibility -------------------------------------
    s4 = S[3]
    set_title(s4)
    set_team(s4)
    drop(shp(s4, "TextBox 8"))
    cw = (W - 2 * 0.38) / 3
    pairs = []
    for i, (hdr, key, items) in enumerate(C.S4_CARDS):
        c = col_card(s4, L + i * (cw + 0.38), TOP, cw, BOT - TOP, key)
        blk = [{"runs": [(hdr, True)], "size": 11.5,
                "color": A if key == "card" else "9A4B2F", "after": 9, "line": 100}]
        for it in items:
            blk.append({"runs": [(it, False)], "size": 10.5, "color": BODY,
                        "bullet": "▪", "after": 8, "line": 104})
        pairs.append((c, blk))
    fit_group(pairs, cap=1.45)
    notes(s4, C.NOTES[4])

    # ---------------- 5 impact ------------------------------------------
    s5 = S[4]
    set_title(s5)
    set_team(s5)
    drop(shp(s5, "TextBox 8"))
    sw = (W - 2 * 0.32) / 3
    for i, (big, lab) in enumerate(C.S5_STATS):
        stat(s5, L + i * (sw + 0.32), TOP, sw, 1.28, big, lab, A, MUTED, C.PAL["card"])
    ct = TOP + 1.28 + 0.30
    ch = BOT - ct
    cw = (W - 0.45) / 2
    pairs = []
    for j, (hdr, rows_) in enumerate([(C.S5_LEFT_HDR, C.S5_LEFT),
                                      (C.S5_RIGHT_HDR, C.S5_RIGHT)]):
        c = col_card(s5, L + j * (cw + 0.45), ct, cw, ch)
        blk = [{"runs": [(hdr, True)], "size": 12, "color": A, "after": 9, "line": 100}]
        for lab, txt in rows_:
            blk.append({"runs": [(lab, True), (txt, False)], "size": 11, "color": BODY,
                        "bullet": "▪", "after": 8, "line": 104})
        pairs.append((c, blk))
    fit_group(pairs, cap=1.45)
    notes(s5, C.NOTES[5])

    # ---------------- 6 references ---------------------------------------
    s6 = S[5]
    set_title(s6)
    set_team(s6)
    drop(shp(s6, "TextBox 8"))
    cw = (W - 0.45) / 2
    pairs = []
    for j, (hdr, rows_) in enumerate([(C.S6_LEFT_HDR, C.S6_LEFT),
                                      (C.S6_RIGHT_HDR, C.S6_RIGHT)]):
        c = col_card(s6, L + j * (cw + 0.45), TOP, cw, BOT - TOP)
        blk = [{"runs": [(hdr, True)], "size": 12, "color": A, "after": 10}]
        for title, url in rows_:
            blk.append({"runs": [(title, False)], "size": 10.5, "color": BODY,
                        "bullet": "▪", "after": 1, "line": 104})
            blk.append({"runs": [(url, False)], "size": 10, "color": A,
                        "indent": 0.22, "after": 8, "line": 104})
        pairs.append((c, blk))
    fit_group(pairs, cap=1.45)
    notes(s6, C.NOTES[6])

    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
