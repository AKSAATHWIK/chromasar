# -*- coding: utf-8 -*-
"""Approximate renderer + overflow report (no LibreOffice available on this box)."""
import sys, os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

DPI = 110
REG = r"C:\Windows\Fonts\arial.ttf"
BLD = r"C:\Windows\Fonts\arialbd.ttf"
_cache = {}


def font(size_pt, bold):
    k = (round(size_pt * 2), bold)
    if k not in _cache:
        px = max(6, int(round(size_pt * DPI / 72.0)))
        _cache[k] = ImageFont.truetype(BLD if bold else REG, px)
    return _cache[k]


def emu_in(v):
    return 0 if v is None else v / 914400.0


def para_props(p):
    pPr = p._p.find(qn("a:pPr"))
    marL = ind = 0.0
    if pPr is not None:
        marL = float(pPr.get("marL", 0)) / 914400.0
        ind = float(pPr.get("indent", 0)) / 914400.0
    def rgb(r):
        try:
            return str(r.font.color.rgb)
        except Exception:
            return "1F497D"
    runs = [(r.text, bool(r.font.bold), (r.font.size.pt if r.font.size else 12.0), rgb(r))
            for r in p.runs]
    size = max([r[2] for r in runs], default=12.0)
    ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
    sb = p.space_before.pt if p.space_before else 0.0
    sa = p.space_after.pt if p.space_after else 0.0
    has_bullet = pPr is not None and pPr.find(qn("a:buChar")) is not None
    return runs, size, ls, sb, sa, marL, ind, has_bullet


def wrap(runs, fnt_w, maxpx):
    """Flow runs into lines of (text,bold,size,color) segments."""
    lines, cur, curw = [], [], 0
    for text, bold, size, color in runs:
      for ln_i, chunk in enumerate(text.split("\n")):
        if ln_i:
            if cur:
                lines.append(cur)
            cur, curw = [], 0
        for i, word in enumerate(chunk.split(" ")):
            if word == "" and i:
                word = " "
            tok = (" " if cur else "") + word
            w = fnt_w(tok, bold, size)
            if curw + w > maxpx and cur:
                lines.append(cur)
                cur, curw = [], 0
                tok = word
                w = fnt_w(tok, bold, size)
            cur.append((tok, bold, size, color))
            curw += w
    if cur:
        lines.append(cur)
    return lines


def render(path, outdir):
    prs = Presentation(path)
    SW, SH = emu_in(prs.slide_width), emu_in(prs.slide_height)
    problems = []
    imgs = []
    for idx, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (int(SW * DPI), int(SH * DPI)), "white")
        d = ImageDraw.Draw(img)

        def fw(t, b, s):
            return d.textlength(t, font=font(s, b))

        for sh in slide.shapes:
            l, t = emu_in(sh.left), emu_in(sh.top)
            w, h = emu_in(sh.width), emu_in(sh.height)
            box = [l * DPI, t * DPI, (l + w) * DPI, (t + h) * DPI]
            st = str(sh.shape_type)
            if "PICTURE" in st:
                d.rectangle(box, fill="#dddddd", outline="#999999")
                d.text((box[0] + 4, box[1] + 4), "IMG", fill="#666666", font=font(9, False))
                continue
            fillhex = None
            try:
                if sh.fill.type is not None and sh.fill.type == 1:
                    fillhex = "#" + str(sh.fill.fore_color.rgb)
            except Exception:
                pass
            if fillhex:
                d.rectangle(box, fill=fillhex)
            elif "AUTO_SHAPE" in st or "FREEFORM" in st:
                d.rectangle(box, outline="#eeeeee")

            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            ml = emu_in(tf.margin_left) or 0.1
            mr = emu_in(tf.margin_right) or 0.1
            mt = emu_in(tf.margin_top) or 0.05
            mb = emu_in(tf.margin_bottom) or 0.05
            avail_w = (w - ml - mr) * DPI
            y = (t + mt) * DPI
            total = 0.0
            for p in tf.paragraphs:
                runs, size, ls, sb, sa, marL, ind, bul = para_props(p)
                if not runs:
                    total += (size * ls + sa) / 72.0
                    y += (size * ls + sa) * DPI / 72.0
                    continue
                indent_px = marL * DPI
                lines = wrap(runs, fw, max(20, avail_w - indent_px))
                y += sb * DPI / 72.0
                total += sb / 72.0
                for li, line in enumerate(lines):
                    x = (l + ml) * DPI + (indent_px + (ind * DPI if li == 0 else 0))
                    if bul and li == 0:
                        d.text((x - 0.16 * DPI, y), "-", fill="#444444", font=font(size, False))
                    align = p.alignment
                    if align is not None and "CENTER" in str(align):
                        lw = sum(fw(seg[0], seg[1], seg[2]) for seg in line)
                        x = (l + w / 2) * DPI - lw / 2
                    for seg, bold, ssz, color in line:
                        d.text((x, y), seg, fill="#" + color, font=font(ssz, bold))
                        x += fw(seg, bold, ssz)
                    y += size * ls * DPI / 72.0
                    total += size * ls / 72.0
                y += sa * DPI / 72.0
                total += sa / 72.0
            avail_h = h - mt - mb
            if total > avail_h + 0.02 and tf.text.strip():
                problems.append("slide %d  %-34s OVERFLOW text %.2f\" > box %.2f\" (by %.2f\")"
                                % (idx, sh.name, total, avail_h, total - avail_h))
            if l < 0.28 or t < -0.1 or l + w > SW - 0.1 or t + h > SH + 0.05:
                problems.append("slide %d  %-34s out of bounds L%.2f T%.2f R%.2f B%.2f"
                                % (idx, sh.name, l, t, l + w, t + h))
            d.rectangle(box, outline="#ff00aa")
        f = os.path.join(outdir, "%s-s%d.png" % (os.path.basename(path)[:24], idx))
        img.save(f)
        imgs.append(f)
    return problems, imgs


if __name__ == "__main__":
    outdir = sys.argv[2] if len(sys.argv) > 2 else "."
    probs, imgs = render(sys.argv[1], outdir)
    print("\n".join(probs) if probs else "no overflow / bounds problems")
    print("\n".join(imgs))
