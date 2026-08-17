"""Shared helpers for filling the SIH 2026 idea template."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

ARIAL = "Arial"


def shp(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    raise KeyError(name + " not on slide")


def _set_bullet(para, char="•", font=ARIAL):
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    if char is None:
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))
        return
    bf = pPr.makeelement(qn("a:buFont"), {"typeface": font})
    pPr.append(bf)
    bc = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(bc)


def _set_space(para, before=0, after=0, line=None):
    if line is not None:
        para.line_spacing = line / 100.0
    para.space_before = Pt(before)
    para.space_after = Pt(after)


def _set_indent(para, mar_in, ind_in):
    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", str(int(mar_in * 914400)))
    pPr.set("indent", str(int(ind_in * 914400)))


def fill(tf, blocks, default_size=13, default_color="1F2A44"):
    """blocks: list of dicts.
      runs: list of (text, bold) or a plain string
      size, color, bullet ('*'|None|char), indent (inches), before, after, line, align
    """
    if not hasattr(tf, "_txBody"):
        tf = tf.text_frame
    tf.word_wrap = True
    # kill autofit so our geometry is authoritative
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    for tag in ("a:spAutoFit", "a:normAutofit"):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    tf.clear()
    # drop every paragraph but the first, then reuse/append
    p0 = tf.paragraphs[0]
    for extra in list(tf._txBody.findall(qn("a:p")))[1:]:
        tf._txBody.remove(extra)
    # scrub inherited props on the survivor
    pPr = p0._p.find(qn("a:pPr"))
    if pPr is not None:
        p0._p.remove(pPr)

    for i, b in enumerate(blocks):
        para = p0 if i == 0 else tf.add_paragraph()
        runs = b.get("runs", b.get("text", ""))
        if isinstance(runs, str):
            runs = [(runs, False)]
        size = b.get("size", default_size)
        color = b.get("color", default_color)
        bullet = b.get("bullet", None)
        indent = b.get("indent", 0.0)
        if bullet:
            _set_indent(para, indent + 0.20, -0.20)
            _set_bullet(para, bullet)
        else:
            _set_indent(para, indent, 0.0)
            _set_bullet(para, None)
        _set_space(para, b.get("before", 0), b.get("after", 3), b.get("line", 100))
        para.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                          "j": PP_ALIGN.JUSTIFY}[b.get("align", "l")]
        for text, bold in runs:
            r = para.add_run()
            r.text = text
            f = r.font
            f.name = ARIAL
            f.size = Pt(size)
            f.bold = bold
            f.underline = b.get("underline", False)
            f.color.rgb = RGBColor.from_string(color)



# ---------------------------------------------------------------- autofit
from PIL import Image, ImageDraw, ImageFont
_MD = ImageDraw.Draw(Image.new("RGB", (8, 8)))
_FC = {}
_DPI = 220.0


def _pf(sz, bold):
    k = (round(sz * 4), bold)
    if k not in _FC:
        _FC[k] = ImageFont.truetype(
            r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
            max(6, int(round(sz * _DPI / 72.0))))
    return _FC[k]


def _tw(t, bold, sz):
    return _MD.textlength(t, font=_pf(sz, bold)) / _DPI  # inches


def _nlines(runs, sz_map, maxw):
    lines, curw, started = 1, 0.0, False
    for text, bold in runs:
        for ci, chunk in enumerate(text.split("\n")):
            if ci:
                lines += 1
                curw, started = 0.0, False
            for word in chunk.split(" "):
                if not word:
                    continue
                tok = (" " if started else "") + word
                w = _tw(tok, bold, sz_map)
                if started and curw + w > maxw:
                    lines += 1
                    curw = _tw(word, bold, sz_map)
                else:
                    curw += w
                started = True
    return lines


def measure(blocks, width_in, ml, mr, scale=1.0):
    total = 0.0
    for b in blocks:
        size = b.get("size", 13) * scale
        ls = b.get("line", 100) / 100.0
        indent = b.get("indent", 0.0) + (0.20 if b.get("bullet") else 0.0)
        runs = b.get("runs", b.get("text", ""))
        if isinstance(runs, str):
            runs = [(runs, False)]
        runs = [(r[0], r[1]) for r in runs]
        n = _nlines(runs, size, max(0.4, width_in - ml - mr - indent))
        total += n * size * ls / 72.0
        total += (b.get("after", 3) + b.get("before", 0)) * scale / 72.0
    return total


def best_scale(shape, blocks, cap=1.60, floor=0.72, target=0.94):
    tf = shape.text_frame
    w = shape.width / 914400.0
    h = shape.height / 914400.0
    ml = (tf.margin_left or 0) / 914400.0 or 0.1
    mr = (tf.margin_right or 0) / 914400.0 or 0.1
    mt = (tf.margin_top or 0) / 914400.0 or 0.05
    mb = (tf.margin_bottom or 0) / 914400.0 or 0.05
    avail = (h - mt - mb) * target
    best = floor
    s = floor
    while s <= cap + 1e-9:
        if measure(blocks, w, ml, mr, s) <= avail:
            best = s
        s += 0.02
    return best


def apply_scale(shape, blocks, scale):
    out = []
    for b in blocks:
        c = dict(b)
        c["size"] = round(b.get("size", 13) * scale, 1)
        c["after"] = round(b.get("after", 3) * scale, 1)
        c["before"] = round(b.get("before", 0) * scale, 1)
        out.append(c)
    fill(shape.text_frame, out)


def fit_fill(shape, blocks, cap=1.60, floor=0.72, target=0.94):
    """Scale font sizes so the block set fills its shape without overflowing."""
    s = best_scale(shape, blocks, cap, floor, target)
    apply_scale(shape, blocks, s)
    return s


def fit_group(pairs, cap=1.60, floor=0.72, target=0.94):
    """Give every (shape, blocks) pair the same font scale - the smallest that fits all."""
    s = min(best_scale(sh, bl, cap, floor, target) for sh, bl in pairs)
    for sh, bl in pairs:
        apply_scale(sh, bl, s)
    return s


def place(shape, l, t, w, h):
    shape.left, shape.top, shape.width, shape.height = (
        Inches(l), Inches(t), Inches(w), Inches(h))


def card(slide, l, t, w, h, fill_hex, line_hex=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(l), Inches(t), Inches(w), Inches(h))
    s.adjustments[0] = 0.06
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    if line_hex:
        s.line.color.rgb = RGBColor.from_string(line_hex)
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    return s


def pipeline(slide, steps, top, box_h=0.95, left=0.5, right=12.83,
             box_hex="0B3C5D", txt_hex="FFFFFF", arrow_hex="9AA7B4", size=9):
    n = len(steps)
    gap = 0.30
    w = (right - left - gap * (n - 1)) / n
    for i, label in enumerate(steps):
        x = left + i * (w + gap)
        b = card(slide, x, top, w, box_h, box_hex)
        b.text_frame.margin_left = Inches(0.04)
        b.text_frame.margin_right = Inches(0.04)
        b.text_frame.margin_top = Inches(0.03)
        b.text_frame.margin_bottom = Inches(0.03)
        b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        fill(b.text_frame, [{"runs": [(label, True)], "size": size,
                             "color": txt_hex, "align": "c", "line": 92}])
        if i < n - 1:
            a = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(x + w + 0.045),
                Inches(top + box_h / 2 - 0.085), Inches(gap - 0.09), Inches(0.17))
            a.fill.solid()
            a.fill.fore_color.rgb = RGBColor.from_string(arrow_hex)
            a.line.fill.background()
            a.shadow.inherit = False


def stat(slide, l, t, w, h, big, label, big_hex, label_hex, bg_hex):
    c = card(slide, l, t, w, h, bg_hex)
    c.text_frame.margin_left = Inches(0.12)
    c.text_frame.margin_right = Inches(0.12)
    c.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    fill(c.text_frame, [
        {"runs": [(big, True)], "size": 26, "color": big_hex, "align": "c",
         "after": 2, "line": 95},
        {"runs": [(label, False)], "size": 11, "color": label_hex, "align": "c",
         "line": 100},
    ])
    return c
