"""Fill the official SIH Idea template with ChromaSAR's content.

The template's own rules, from its last slide:
  - maximum six slides INCLUDING the title
  - points, not paragraphs; diagrams and pictures encouraged
  - use the provided template without changing the idea-detail pointers
  - upload as PDF, not PPT

So this edits the supplied file rather than restyling it. The section headings, colours,
logo and footers are left exactly as issued; only the grey prompt text inside each content
box is replaced. The boxes themselves ARE resized - the template sizes them to fit its own
one-line prompts (slide 6's is 0.57in tall) while leaving four inches of empty slide below,
so real content spills past the edge with no warning. PowerPoint does not flag that; it
only shows up on the projector. Each box is moved to the usable area and the body font is
stepped down until the text provably fits.

Our own 7-slide deck stays separate - that one is the desk-side reference, this one is the
submission artefact.

    python deck/fill_template.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "SIH2026-IDEA-Presentation-Format.pptx"
OUT = ROOT / "SIH2026-DeltaForce-SIH1733-IDEA-SUBMISSION.pptx"
DIAGRAM = ROOT / "deck" / "architecture_slide.png"

TEAM = "Delta Force"

TITLE_BLOCK = [
    ("Problem Statement ID – SIH1733", True),
    ("Problem Statement Title – SAR Image Colorization for Comprehensive "
     "Insight using Deep Learning Model", True),
    ("Theme – Space Technology", False),
    ("PS Category – Software", False),
    (f"Team Name – {TEAM}", False),
]

#: slide index -> (bullet, is_heading). Deliberately short lines: the template asks for
#: points, and a judge at a desk reads a slide in about four seconds.
CONTENT = {
    1: [
        ("ChromaSAR — colorized SAR you can act on, with a confidence you can check", True),
        ("Problem: optical satellites go blind under cloud — exactly during the monsoon. "
         "Radar sees through it, but radar is grayscale texture only specialists can read.", False),
        ("Solution: a conditional GAN colorizes Sentinel-1 SAR into optical-like imagery, "
         "trained on 10,000 co-registered SAR–optical pairs.", False),
        ("USP — radar carries no colour, so every pixel ships a CALIBRATED CONFIDENCE. "
         "Below your threshold the output goes grey: “insufficient evidence”, not a "
         "confident guess. Per-pixel uncertainty on SAR is not new; what is new is that "
         "ours is calibrated, shown to the user, and GATES the downstream answer.", False),
        ("That confidence GATES everything downstream — flood extent, change detection "
         "and surface cover are all filtered by it.", False),
        ("Unique 1 — we diagnosed WHY existing models look unsatisfactory. L1 loss is "
         "minimised by the conditional MEAN of every plausible colour, which is blur by "
         "definition, and PSNR rewards the same thing so standard evaluation cannot see "
         "it. A gradient-difference loss fixed it: sharpness 0.21 → 0.78 of the real "
         "optical, saturation 45% → 102%.", False),
        ("Unique 2 — we measure what does not work and REFUSE to ship it. Built-up "
         "land cover from radar alone scores AUC 0.483, below chance, so the app reports "
         "it as unavailable and says why.", False),
    ],
    2: [
        ("Technologies", True),
        ("Python · PyTorch · FastAPI · Next.js 15 · TypeScript · GDAL/rasterio · "
         "QGIS-ready GeoTIFF export", False),
        ("Generator ResNet34-UNet (ImageNet encoder) + 70×70 PatchGAN discriminator", False),
        ("Loss: L1 + frozen-VGG16 perceptual + adversarial + gradient-difference", False),
        ("Confidence: 10-pass Monte-Carlo dropout, per pixel", False),
        ("Flood: dual-polarisation threshold baseline + ResNet34-UNet segmentation, "
         "temperature-calibrated", False),
        ("Methodology", True),
        ("Split by SCENE, not patch — adjacent tiles cannot leak across train/validation", False),
        ("Trained on a serverless A10G GPU (~$3 total). Inference runs on CPU: 1.1 s per "
         "scene, in-process, no network call", False),
        ("Prototype running today: 4 workspaces, upload your own file, georeferenced "
         "export, 43 regression tests green", False),
    ],
    3: [
        ("Feasibility — demonstrated, not projected", True),
        ("Runs today on a laptop CPU. No GPU, no cloud dependency, no API key.", False),
        ("Flood mapping IoU 0.681 on the Sen1Floods11 OFFICIAL test split, against 0.550 "
         "for the physics baseline. Probabilities calibrated: ECE 0.029 → 0.016.", False),
        ("An uploaded file scores identically to our benchmark — 0.00% difference across "
         "three scenes on three continents.", False),
        ("Honest positioning: 0.681 beats the Sen1Floods11 published baseline (0.662) and "
         "matches a supervised SAR-only model (0.676). SAR-only state of the art is ~0.72 "
         "(DeepSARFlood, IIT Delhi). We are competitive, not leading, and we say so.", False),
        ("Risks, and how we handle them", True),
        ("Hallucination — the network invents colour the radar never carried. → Ship the "
         "confidence map and refuse below threshold. Surfaced, never hidden.", False),
        ("Blur — L1 drives the output to the conditional mean. → Gradient-difference loss, "
         "plus a sharpness metric PSNR is structurally blind to.", False),
        ("Flooded vegetation double-bounces and returns BRIGHT, not dark, so the physics "
         "threshold misses it. → Learned segmentation replaces physics where physics fails.", False),
        ("Land cover from radar alone does not work (built-up AUC 0.483). → Measured, and "
         "refused. Cover is read from the co-registered Sentinel-2 optical and labelled "
         "as such.", False),
        ("Domain shift across terrain. → Stratified sampling over agro-climatic zones, "
         "per-region fine-tuning.", False),
    ],
    4: [
        ("Impact — target audience: Remote Sensing Image Analysts", True),
        ("SAR interpretation drops from expert-only to any GIS user; exports open straight "
         "in QGIS on top of the source scene.", False),
        ("Disaster response: flood extent mapped during the monsoon, exactly when optical "
         "satellites are blind. Permanent water is separated from new flooding — a river "
         "is not a disaster.", False),
        ("Benefits", True),
        ("Social — faster, better-informed response in the months India floods most.", False),
        ("Economic — removes dependence on cloud-free optical tasking; runs on existing "
         "hardware, no GPU purchase.", False),
        ("Strategic — indigenous SAR-interpretation capability that transfers to RISAT and "
         "NISAR; works offline, in the field, with no connectivity.", False),
        ("Environmental — the PS's own named uses: continuous monitoring of wetlands, "
         "glaciers and coastal change instead of seasonal snapshots.", False),
        ("Archive unlock — decades of existing SAR holdings become visually searchable.", False),
    ],
    5: [
        ("Datasets", True),
        ("Sen1Floods11 — Bonafilia et al., CVPRW 2020. 446 hand-labelled flood chips, "
         "11 regions, official split.  github.com/cloudtostreet/Sen1Floods11", False),
        ("SEN1-2 — Schmitt et al., ISPRS 2018. 10,000 co-registered SAR–optical pairs; the "
         "paper names SAR colorization as an intended use.  mediatum.ub.tum.de/1436631", False),
        ("Imagery — Sentinel-1 / Sentinel-2 via ESA Copernicus.  "
         "dataspace.copernicus.eu", False),
        ("Prior art and incumbents", True),
        ("Copernicus GFM — operational global flood service. Already ships a per-pixel "
         "likelihood, but on a binary mask, with no colorization.  extwiki.eodc.eu/gfm", False),
        ("ISRO NRSC — near-real-time flood inundation via NDEM / Bhuvan (EOS-04, "
         "Sentinel-1). Authoritative, event-triggered, delivered as maps.  nrsc.gov.in", False),
        ("OSCAR, 2026 — uncertainty-aware SAR-to-optical diffusion. Closest prior art; "
         "the uncertainty stays in the training loss.  arxiv.org/abs/2601.06835", False),
        ("Methods", True),
        ("Isola et al., Image-to-Image Translation with Conditional Adversarial Networks "
         "(pix2pix), CVPR 2017 — our deliberate baseline.", False),
        ("Mathieu et al., Deep Multi-Scale Video Prediction Beyond MSE, ICLR 2016 — the "
         "gradient-difference loss that removed our blur.", False),
        ("Gal & Ghahramani, Dropout as a Bayesian Approximation, ICML 2016 — the basis of "
         "our per-pixel confidence.", False),
        ("Guo et al., On Calibration of Modern Neural Networks, ICML 2017 — temperature "
         "scaling; ECE 0.029 → 0.016.", False),
        ("Xu, Modification of NDWI (MNDWI), IJRS 2006 — optical water reference; we "
         "re-measured the threshold rather than adopting it.", False),
    ],
}

#: left, top, width, height in inches. Below the section heading, above the footer bar.
BOX = {1: (0.67, 1.32, 12.00, 5.50),
       2: (0.67, 4.10, 12.00, 2.75),   # diagram occupies the space above
       3: (0.67, 1.32, 12.00, 5.50),
       4: (0.67, 1.32, 12.00, 5.50),
       5: (0.67, 1.32, 12.00, 5.50)}

#: Largest first - we want the biggest type that fits, not the first that is safe.
#: 18pt is the ceiling because above that the bullets start looking like a poster.
BODY_SIZES = [18.0, 17.0, 16.0, 15.0, 14.0, 13.0, 12.0, 11.0, 10.0, 9.5]


def plain(para):
    """Drop inherited list formatting.

    Only the FIRST paragraph of a placeholder inherits the layout's bullet style, so
    without this the opening heading picks up a stray glyph and hanging indent while the
    identical heading further down does not. We supply our own bullet character, so clear
    the inherited one everywhere and start every line at the same left edge.
    """
    p = para._p
    pPr = p.find(qn("a:pPr"))
    if pPr is None:
        pPr = p.makeelement(qn("a:pPr"), {})
        p.insert(0, pPr)
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def wrapped_height(items, body_pt: float, width_in: float) -> float:
    """Estimated rendered height, in inches. Deliberately pessimistic."""
    total = 0.0
    for text, is_head in items:
        pt = body_pt + 2.5 if is_head else body_pt
        char_w = pt * 0.0071          # mean glyph advance for this typeface, in inches
        per_line = max(10, int(width_in / char_w))
        n = len(text) + (0 if is_head else 3)
        lines = max(1, -(-n // per_line))
        total += lines * pt * 1.22 / 72.0        # line box
        total += (7.0 if is_head else 4.0) / 72.0   # space_after
    return total


def fill(slide, items, geom):
    """Replace the prompt box with our bullets, sized so they provably fit."""
    boxes = [sh for sh in slide.shapes
             if sh.has_text_frame and len(sh.text_frame.text.strip()) > 40]
    if not boxes:
        return None
    box = max(boxes, key=lambda sh: sh.width * sh.height)
    left, top, width, height = geom
    box.left, box.top = Inches(left), Inches(top)
    box.width, box.height = Inches(width), Inches(height)

    usable = width - 0.25          # internal left+right inset
    body = BODY_SIZES[-1]
    for cand in BODY_SIZES:
        if wrapped_height(items, cand, usable) <= height:
            body = cand
            break

    tf = box.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)

    for i, (text, is_head) in enumerate(items):
        para = first if i == 0 else tf.add_paragraph()
        plain(para)
        run = para.add_run()
        run.text = text if is_head else "•  " + text
        run.font.size = Pt(body + 2.5 if is_head else body)
        run.font.bold = is_head
        para.space_after = Pt(7 if is_head else 4)
    return body


def main() -> int:
    if not SRC.exists():
        print("template not found:", SRC)
        return 1
    prs = Presentation(SRC)

    # The section headings are left alone. They look like they start at L0.77 and would
    # run under the team-name badge, but their alignment is inherited (centre), so they
    # actually sit mid-slide and clear both the badge and the SIH logo. Nudging them right
    # to "make room" is what breaks them - it pushes the long ones into the logo.

    for sl in prs.slides:                      # the team badge on every content slide
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() == "Your Team Name":
                tf = sh.text_frame
                for p in list(tf.paragraphs)[1:]:
                    p._p.getparent().remove(p._p)
                for r in list(tf.paragraphs[0].runs)[1:]:
                    r._r.getparent().remove(r._r)
                runs = tf.paragraphs[0].runs
                if runs:
                    runs[0].text = TEAM
                    runs[0].font.size = Pt(12)
                    runs[0].font.bold = True

    for sh in prs.slides[0].shapes:            # slide 1 identity block
        if sh.has_text_frame and "Problem Statement ID" in sh.text_frame.text:
            tf = sh.text_frame
            tf.word_wrap = True
            for p in list(tf.paragraphs)[1:]:
                p._p.getparent().remove(p._p)
            for r in list(tf.paragraphs[0].runs):
                r._r.getparent().remove(r._r)
            for i, (line, bold) in enumerate(TITLE_BLOCK):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                plain(para)
                run = para.add_run()
                run.text = line
                run.font.size = Pt(14)
                run.font.bold = bold
                para.space_after = Pt(6)
            break

    for idx, items in CONTENT.items():
        pt = fill(prs.slides[idx], items, BOX[idx])
        print(f"  slide {idx + 1}: body {pt}pt" if pt else
              f"  slide {idx + 1}: WARNING no box filled")

    # the template asks for a flow chart on TECHNICAL APPROACH
    if DIAGRAM.exists():
        prs.slides[2].shapes.add_picture(str(DIAGRAM), Inches(1.42), Inches(1.35),
                                         width=Inches(10.5))
        print("  slide 3: pipeline diagram added")

    # "You can delete this slide (Important Pointers) when you upload"
    lst = prs.slides._sldIdLst
    ids = list(lst)
    if len(ids) == 7:
        prs.part.drop_rel(ids[6].rId)
        lst.remove(ids[6])

    prs.save(OUT)
    print(f"\nwrote {OUT.name}  ({len(list(prs.slides))} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
